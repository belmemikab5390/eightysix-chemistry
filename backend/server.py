"""
EightySix Chemistry - Production Server
Cloud-ready with R2 storage integration
"""

from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import requests
import json
import os
import re
from datetime import datetime
from pptx import Presentation

app = Flask(__name__)
CORS(app)

# ============================================
# CONFIGURATION - ENVIRONMENT VARIABLES
# ============================================

# Get from environment (set in Render dashboard)
OPENROUTER_API_KEY = os.environ.get('OPENROUTER_API_KEY', 'your-key-here')
PORT = int(os.environ.get('PORT', 5000))
PRODUCTION = os.environ.get('PRODUCTION', 'false').lower() == 'true'

# OpenRouter configuration
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"
MODEL = os.environ.get('MODEL', 'xiaomi/mimo-v2-flash:free')

# ============================================
# BOOK LIBRARY - R2 URLs (NO LOCAL PATHS!)
# ============================================

BOOK_LIBRARY = {
    'zumdahl': {
        'name': 'General Chemistry',
        'author': 'Zumdahl & Zumdahl',
        'chunks_url': 'https://drive.google.com/uc?id=1jYCaCjhM5Po6ucgzNbAW7n9DVE6wTr0G',
        'pdf_url': 'https://drive.google.com/uc?export=view&id=1ElyPAg8BLiycLWtbZ7nCXJ9k9tqEAu3n'
    },
    'atkins': {
        'name': 'Physical Chemistry',
        'author': 'Atkins & de Paula',
        'chunks_url': 'https://drive.google.com/uc?id=1jYCaCjhM5Po6ucgzNbAW7n9DVE6wTr0G',
        'pdf_url': 'https://drive.google.com/uc?export=view&id=1ElyPAg8BLiycLWtbZ7nCXJ9k9tqEAu3n'
    },
    'harris': {
        'name': 'Quantitative Chemical Analysis',
        'author': 'Daniel C. Harris',
        'chunks_url': 'https://drive.google.com/uc?id=1oSSwyWZSvMNEEvkunCE4U_xO4h3uh4pF',
        'pdf_url': 'https://drive.google.com/uc?export=view&id=1w9vWa_T76YmOe-1OoCSkOYhlfbtWk_UM'
    },
    'klein': {
        'name': 'Organic Chemistry',
        'author': 'David Klein',
        'chunks_url': 'https://drive.google.com/uc?id=1KdiO7gnP26-1M5_nY55hGQCjqLI11ALE',
        'pdf_url': 'https://drive.google.com/uc?export=view&id=1FXaL5Xt_8kraeWisR22HGkNqDdNcF4ZQ'
    }
}

# Currently loaded book
current_book_info = {
    'id': 'zumdahl',
    'name': 'General Chemistry',
    'author': 'Zumdahl & Zumdahl'
}

# ============================================
# AI TEXTBOOK SEARCH
# ============================================

class AITextbookSearch:
    HIGH_CONFIDENCE = 0.65
    LOW_CONFIDENCE = 0.35
    
    def __init__(self):
        print("🔧 Initializing AI Textbook Search...")
        self.model = None  # no local model
        print("✅ Running in API mode (no local model)")
        
        self.chunks = []
        self.textbook = None
        print("⚠️  No chunks loaded - will load on book selection")
    
    def load_chunks_from_url(self, url):
        try:
            print(f"📥 Fetching chunks from: {url}")

            response = requests.get(url, timeout=60)
            response.raise_for_status()

            content_type = response.headers.get("Content-Type", "")

            chunks = json.loads(response.text)

            chunks = response.json()

            self.chunks = chunks
            self.textbook = {'pages': chunks}

            print(f"✅ Loaded {len(chunks)} chunks")
            return True

        except Exception as e:
            print("LOAD ERROR:", e)
            return False
    
    def smart_search(self, question):
        if not self.chunks:
            return "No textbook loaded. Please select a book first.", 0.0, False

        question_words = set(question.lower().split())

        scored = []
        for chunk in self.chunks:
            text_words = set(chunk["text"].lower().split())
            overlap = len(question_words & text_words)
            score = overlap / (len(question_words) + 1)
            scored.append((chunk, score))

        scored.sort(key=lambda x: x[1], reverse=True)

        if not scored:
            return "No relevant content found.", 0.0, False

        top_chunk, top_score = scored[0]
        is_relevant = top_score > 0.1

        context = "\n\n".join([
            f"[Page {c['page']}] {c['text']}"
            for c, _ in scored[:3]
        ])

        return context, top_score, is_relevant
    
        def get_candidate_pages(self, topic, top_k=5):
            topic_words = set(topic.lower().split())

            scored = []
            for chunk in self.chunks:
                text_words = set(chunk["text"].lower().split())
                overlap = len(topic_words & text_words)
                score = overlap / (len(topic_words) + 1)

                scored.append({
                    'page': chunk['page'],
                    'text': chunk['text'],
                    'score': score
                })

            scored.sort(key=lambda x: x['score'], reverse=True)
            return scored[:top_k]


# Initialize AI search
ai_search = AITextbookSearch()


# ============================================
# HELPER FUNCTIONS
# ============================================

def call_ai(prompt, system_prompt="You are an expert chemistry tutor."):
    """Call OpenRouter API"""
    try:
        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json",
            "HTTP-Referer": "https://chemistry-app.com",
            "X-Title": "EightySix Chemistry"
        }
        
        payload = {
            "model": MODEL,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 3000
        }
        
        response = requests.post(
            OPENROUTER_URL,
            headers=headers,
            json=payload,
            timeout=60
        )
        
        if response.status_code == 200:
            return response.json()['choices'][0]['message']['content']
        else:
            return f"Error: API returned {response.status_code}"
            
    except Exception as e:
        return f"Error: {str(e)}"


# ============================================
# ROUTES
# ============================================

@app.route('/')
def home():
    """Serve info page"""
    return jsonify({
        'name': 'EightySix Chemistry API',
        'version': '1.0',
        'status': 'running',
        'endpoints': {
            'health': '/health',
            'ask': '/ask',
            'load_book': '/load-book',
            'library': '/get-library',
            'flashcards': '/generate-flashcards',
            'ppt_upload': '/upload-ppt',
            'materials': '/generate-study-materials'
        }
    })


@app.route('/health', methods=['GET'])
def health():
    """Health check"""
    return jsonify({
        'status': 'healthy',
        'mode': 'production' if PRODUCTION else 'development',
        'textbook_loaded': ai_search.textbook is not None,
        'chunks_count': len(ai_search.chunks),
        'current_book': current_book_info,
        'storage_configured': True,
        'api_configured': OPENROUTER_API_KEY != 'your-key-here'
    })


@app.route('/load-book', methods=['POST'])
def load_book():
    """Load a book from R2"""
    global current_book_info
    
    try:
        data = request.json
        book_id = data.get('bookId')
        
        print(f"\n📚 Request to load book: {book_id}")
        
        if book_id not in BOOK_LIBRARY:
            return jsonify({
                'success': False,
                'error': f'Book "{book_id}" not found'
            }), 404
        
        book = BOOK_LIBRARY[book_id]
        chunks_url = book['chunks_url']
        
        print(f"📖 Loading: {book['name']} by {book['author']}")
        print(f"📂 From R2: {chunks_url}")
        
        # Load chunks from R2
        success = ai_search.load_chunks_from_url(chunks_url)
        
        if not success:
            return jsonify({
                'success': False,
                'error': 'Failed to load chunks from R2'
            }), 500
        
        # Update current book
        current_book_info = {
            'id': book_id,
            'name': book['name'],
            'author': book['author']
        }
        
        print(f"✅ Successfully loaded {len(ai_search.chunks)} chunks")
        
        return jsonify({
            'success': True,
            'book_id': book_id,
            'book_name': book['name'],
            'author': book['author'],
            'chunks_count': len(ai_search.chunks)
        })
        
    except Exception as e:
        print(f"❌ Error loading book: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500


@app.route('/get-library', methods=['GET'])
def get_library():
    """Get list of available books"""
    books = []
    for book_id, book_info in BOOK_LIBRARY.items():
        books.append({
            'id': book_id,
            'name': book_info['name'],
            'author': book_info['author'],
            'available': True  # All books available in R2
        })
    
    return jsonify({
        'success': True,
        'books': books,
        'current_book': current_book_info
    })


@app.route('/ask', methods=['POST'])
def ask():
    """Answer questions"""
    try:
        data = request.json
        question = data.get('question', '')
        complexity = data.get('complexity', 3)
        
        print(f"\n📝 Question: {question}")
        
        if not ai_search.chunks:
            return jsonify({
                'success': False,
                'error': 'No textbook loaded. Please select a book first.'
            }), 400
        
        # Search for answer
        context, similarity, is_relevant = ai_search.smart_search(question)
        
        print(f"🔍 Similarity: {similarity:.3f}, Relevant: {is_relevant}")
        
        # Generate response
        complexity_levels = {
            1: "Explain simply, like to a beginner.",
            2: "Explain clearly for high school level.",
            3: "Balanced explanation with proper terminology.",
            4: "Detailed explanation for advanced students.",
            5: "Comprehensive university-level explanation."
        }
        
        instruction = complexity_levels.get(complexity, complexity_levels[3])
        
        if is_relevant:
            prompt = f"""You are a chemistry tutor.

TEXTBOOK CONTEXT:
{context}

STUDENT QUESTION: {question}

INSTRUCTIONS: {instruction}

Provide a clear, accurate answer based on the textbook context."""
        else:
            prompt = f"""The student asked: "{question}"

This doesn't seem related to the current chemistry textbook. 
Politely let them know and offer to help with chemistry topics from the book."""
        
        answer = call_ai(prompt)
        
        return jsonify({
            'success': True,
            'answer': answer,
            'context': context,
            'similarity': float(similarity),
            'is_relevant': is_relevant
        })
        
    except Exception as e:
        print(f"❌ Error: {e}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500


@app.route('/generate-flashcards', methods=['POST'])
def generate_flashcards():
    """Generate flashcards"""
    try:
        data = request.json
        topic = data.get('topic', '')
        count = data.get('count', 5)
        
        print(f"\n🎴 Generating {count} flashcards for: {topic}")
        
        if not ai_search.chunks:
            return jsonify({
                'success': False,
                'error': 'No textbook loaded'
            }), 400
        
        # Get relevant pages
        pages = ai_search.get_candidate_pages(topic, top_k=3)
        
        if not pages or pages[0]['score'] < 0.3:
            return jsonify({
                'success': False,
                'error': f'Topic "{topic}" not found in textbook'
            }), 404
        
        # Combine content
        content = "\n\n".join([p['text'][:1000] for p in pages])
        
        prompt = f"""Create {count} flashcards about {topic} from this content:

{content}

Format EXACTLY as:
Q: [Question]
A: [Answer]

Q: [Question]
A: [Answer]"""
        
        response = call_ai(prompt)
        
        # Parse flashcards
        flashcards = []
        lines = response.split('\n')
        current_q = None
        
        for line in lines:
            line = line.strip()
            if line.startswith('Q:'):
                current_q = line[2:].strip()
            elif line.startswith('A:') and current_q:
                flashcards.append({
                    'front': current_q,
                    'back': line[2:].strip()
                })
                current_q = None
        
        print(f"✅ Generated {len(flashcards)} flashcards")
        
        return jsonify({
            'success': True,
            'flashcards': flashcards[:count]
        })
        
    except Exception as e:
        print(f"❌ Error: {e}")
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500


@app.route('/upload-ppt', methods=['POST'])
def upload_ppt():
    """Process PowerPoint"""
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file'}), 400
        
        file = request.files['file']
        
        if not file.filename.endswith(('.ppt', '.pptx')):
            return jsonify({'success': False, 'error': 'Must be .ppt/.pptx'}), 400
        
        print(f"\n📊 Processing: {file.filename}")
        
        # Save temp
        temp_path = f"temp_{datetime.now().strftime('%Y%m%d%H%M%S')}_{file.filename}"
        file.save(temp_path)
        
        try:
            prs = Presentation(temp_path)
            
            slides_content = []
            for i, slide in enumerate(prs.slides, 1):
                slide_data = {
                    'slide_number': i,
                    'title': '',
                    'content': [],
                    'notes': ''
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        if hasattr(shape, 'placeholder_format'):
                            try:
                                if shape.placeholder_format.type == 1:
                                    slide_data['title'] = text
                                    continue
                            except:
                                pass
                        
                        slide_data['content'].append(text)
                
                try:
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame
                        if notes and notes.text:
                            slide_data['notes'] = notes.text.strip()
                except:
                    pass
                
                slides_content.append(slide_data)
            
            print(f"✅ Extracted {len(slides_content)} slides")
            
            # Clean up
            os.remove(temp_path)
            
            return jsonify({
                'success': True,
                'slides': slides_content,
                'total_slides': len(slides_content),
                'filename': file.filename
            })
            
        except Exception as e:
            if os.path.exists(temp_path):
                os.remove(temp_path)
            raise e
        
    except Exception as e:
        print(f"❌ Error: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/generate-study-materials', methods=['POST'])
def generate_study_materials():
    """Generate study materials from slides"""
    try:
        data = request.json
        slides = data.get('slides', [])
        material_type = data.get('type', 'notes')
        
        print(f"\n📚 Generating {material_type} from {len(slides)} slides")
        
        # Combine content
        full_content = ""
        for slide in slides:
            full_content += f"\n\n=== Slide {slide['slide_number']}"
            if slide['title']:
                full_content += f": {slide['title']}"
            full_content += " ===\n"
            full_content += "\n".join(slide['content'])
            if slide['notes']:
                full_content += f"\n\nNotes: {slide['notes']}"
        
        results = {}
        
        # Generate based on type
        if material_type in ['notes', 'all']:
            prompt = f"Create study notes from this lecture:\n\n{full_content}\n\nUse clear headings and formatting."
            results['notes'] = call_ai(prompt)
        
        if material_type in ['reviewer', 'all']:
            prompt = f"Create an exam reviewer from:\n\n{full_content}\n\nFocus on key concepts and practice problems."
            results['reviewer'] = call_ai(prompt)
        
        if material_type in ['flashcards', 'all']:
            prompt = f"Create 15 flashcards from:\n\n{full_content}\n\nFormat: Q: ... A: ..."
            text = call_ai(prompt)
            
            flashcards = []
            lines = text.split('\n')
            current_q = None
            
            for line in lines:
                line = line.strip()
                if line.startswith('Q:'):
                    current_q = line[2:].strip()
                elif line.startswith('A:') and current_q:
                    flashcards.append({
                        'front': current_q,
                        'back': line[2:].strip()
                    })
                    current_q = None
            
            results['flashcards'] = flashcards
        
        if material_type in ['summary', 'all']:
            prompt = f"Create a one-page summary:\n\n{full_content}"
            results['summary'] = call_ai(prompt)
        
        if material_type in ['quiz', 'all']:
            prompt = f"Create 10 MCQ questions:\n\n{full_content}"
            results['quiz'] = call_ai(prompt)
        
        return jsonify({
            'success': True,
            'materials': results
        })
        
    except Exception as e:
        print(f"❌ Error: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500


# ============================================
# START SERVER
# ============================================

if __name__ == '__main__':
    print("=" * 60)
    print("🤖 EIGHTYSIX CHEMISTRY - PRODUCTION SERVER")
    print("=" * 60)
    print(f"🌐 Mode: {'PRODUCTION' if PRODUCTION else 'DEVELOPMENT'}")
    print(f"📚 Books: {len(BOOK_LIBRARY)}")
    print(f"✅ Storage: {'Configured' if STORAGE_URL else 'Not configured'}")
    print(f"✅ API: {'Configured' if OPENROUTER_API_KEY != 'your-key-here' else 'Not configured'}")
    print(f"🌐 Port: {PORT}")
    print("=" * 60)
    
    app.run(host='0.0.0.0', port=PORT)
